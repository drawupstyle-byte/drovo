from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from docx import Document

# === ЧТЕНИЕ ИСХОДНОГО ФАЙЛА ===
doc = Document('.github/workflows/drop.docx')
full_text = []
for para in doc.paragraphs:
    if para.text.strip():
        full_text.append(para.text)

# === НАСТРОЙКИ ПРЕЗЕНТАЦИИ ===
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 формат
prs.slide_height = Inches(7.5)

# === ЦВЕТОВАЯ ПАЛИТРА (по ТЗ) ===
COLOR_BG_LIGHT = RGBColor(245, 245, 245)      # #F5F5F5 - светло-серый фон
COLOR_TEXT_DARK = RGBColor(51, 51, 51)         # #333333 - графитовый текст
COLOR_TEXT_BLACK = RGBColor(0, 0, 0)           # #000000 - черный
COLOR_ACCENT = RGBColor(52, 191, 163)          # #34BFA3 - бирюзовый акцент
COLOR_DIVIDER = RGBColor(224, 224, 224)        # #E0E0E0 - разделитель
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_PRIMARY_DARK = RGBColor(20, 24, 50)      # Темный для титульного
COLOR_PLACEHOLDER = RGBColor(240, 240, 240)    # Светло-серая плашка под текстом

# === ШРИФТЫ (по ТЗ) ===
FONT_NAME = 'Google Sans'
FONT_FALLBACK = 'Arial'  # На случай если Google Sans не установлен

TITLE_SIZE = Pt(56)
SUBTITLE_SIZE = Pt(36)
HEADING_SIZE = Pt(48)
SUBHEADING_SIZE = Pt(36)
BODY_SIZE = Pt(24)
SMALL_SIZE = Pt(16)

def add_background(slide, color=COLOR_BG_LIGHT):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg

def add_placeholder_box(slide, left, top, width, height):
    """Добавляет светло-серую плашку под текст"""
    placeholder = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = COLOR_PLACEHOLDER
    placeholder.line.fill.background()
    # Скругление углов
    placeholder.adjustments[0] = 0.1
    return placeholder

def set_font_for_run(run, size=BODY_SIZE, bold=False, color=COLOR_TEXT_DARK, font_name=FONT_FALLBACK):
    run.font.name = font_name
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    return run

def create_title_slide(prs, title_text, subtitle_lines):
    """Титульный слайд с темным фоном"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Темный фон
    add_background(slide, COLOR_PRIMARY_DARK)
    
    # Бирюзовый акцент слева (вертикальная полоса)
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.4), prs.slide_height)
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_ACCENT
    accent_bar.line.fill.background()
    
    # Заголовок по центру
    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.5), Inches(2))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_para = title_frame.paragraphs[0]
    title_para.text = title_text
    title_para.alignment = PP_ALIGN.CENTER
    
    run = title_para.runs[0]
    run.font.name = FONT_FALLBACK
    run.font.size = TITLE_SIZE
    run.font.bold = True
    run.font.color.rgb = COLOR_WHITE
    
    # Подзаголовок
    sub_y = Inches(4.3)
    for line in subtitle_lines:
        if line.strip():
            sub_box = slide.shapes.add_textbox(Inches(1.5), sub_y, Inches(10.5), Inches(0.6))
            sub_frame = sub_box.text_frame
            sub_frame.clear()
            sub_para = sub_frame.paragraphs[0]
            sub_para.text = line
            sub_para.alignment = PP_ALIGN.CENTER
            
            run = sub_para.runs[0]
            run.font.name = FONT_FALLBACK
            run.font.size = SUBTITLE_SIZE
            run.font.bold = False
            run.font.color.rgb = COLOR_WHITE
            
            sub_y += Inches(0.5)
    
    return slide

def create_section_slide(prs, section_num, section_title, content_blocks):
    """Слайд с контентом на светлом фоне с плашками"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Светлый фон
    add_background(slide, COLOR_BG_LIGHT)
    
    # Бирюзовый акцент слева
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.4), prs.slide_height)
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_ACCENT
    accent_bar.line.fill.background()
    
    # Заголовок секции
    header_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.4), Inches(11.5), Inches(1))
    header_frame = header_box.text_frame
    header_frame.clear()
    header_para = header_frame.paragraphs[0]
    header_para.text = f"{section_num} / {section_title}"
    
    run = header_para.runs[0]
    run.font.name = FONT_FALLBACK
    run.font.size = HEADING_SIZE
    run.font.bold = True
    run.font.color.rgb = COLOR_TEXT_DARK
    
    # Разделитель
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.3), Inches(11.5), Inches(0.05))
    divider.fill.solid()
    divider.fill.fore_color.rgb = COLOR_ACCENT
    divider.line.fill.background()
    
    # Контент с плашками
    y_pos = Inches(1.8)
    for block in content_blocks:
        block_type = block.get('type', 'body')
        text = block.get('text', '')
        
        if block_type == 'subheading':
            # Плашка для подзаголовка
            placeholder = add_placeholder_box(slide, Inches(1.0), y_pos, Inches(11.5), Inches(0.7))
            box = slide.shapes.add_textbox(Inches(1.2), y_pos + Inches(0.15), Inches(11.1), Inches(0.4))
            frame = box.text_frame
            frame.clear()
            para = frame.paragraphs[0]
            para.text = text
            
            run = para.runs[0]
            run.font.name = FONT_FALLBACK
            run.font.size = SUBHEADING_SIZE
            run.font.bold = True
            run.font.color.rgb = COLOR_TEXT_DARK
            
            y_pos += Inches(0.85)
            
        elif block_type == 'body':
            # Плашка для основного текста
            text_height = max(Inches(0.6), Inches(0.3) + Inches(0.25) * (len(text) // 100))
            placeholder = add_placeholder_box(slide, Inches(1.0), y_pos, Inches(11.5), text_height)
            box = slide.shapes.add_textbox(Inches(1.2), y_pos + Inches(0.15), Inches(11.1), text_height - Inches(0.3))
            frame = box.text_frame
            frame.word_wrap = True
            frame.clear()
            para = frame.paragraphs[0]
            para.text = text
            
            run = para.runs[0]
            run.font.name = FONT_FALLBACK
            run.font.size = BODY_SIZE
            run.font.bold = False
            run.font.color.rgb = COLOR_TEXT_DARK
            
            y_pos += text_height + Inches(0.15)
            
        elif block_type == 'bullet':
            # Плашка для маркированного списка
            text_height = max(Inches(0.5), Inches(0.25) + Inches(0.25) * (len(text) // 80))
            placeholder = add_placeholder_box(slide, Inches(1.0), y_pos, Inches(11.5), text_height)
            box = slide.shapes.add_textbox(Inches(1.2), y_pos + Inches(0.15), Inches(11.1), text_height - Inches(0.3))
            frame = box.text_frame
            frame.word_wrap = True
            frame.clear()
            para = frame.paragraphs[0]
            para.text = '• ' + text
            
            run = para.runs[0]
            run.font.name = FONT_FALLBACK
            run.font.size = BODY_SIZE
            run.font.bold = False
            run.font.color.rgb = COLOR_TEXT_DARK
            
            y_pos += text_height + Inches(0.15)
            
        elif block_type == 'small':
            # Плашка для мелкого текста
            placeholder = add_placeholder_box(slide, Inches(1.0), y_pos, Inches(11.5), Inches(0.5))
            box = slide.shapes.add_textbox(Inches(1.2), y_pos + Inches(0.15), Inches(11.1), Inches(0.2))
            frame = box.text_frame
            frame.word_wrap = True
            frame.clear()
            para = frame.paragraphs[0]
            para.text = text
            
            run = para.runs[0]
            run.font.name = FONT_FALLBACK
            run.font.size = SMALL_SIZE
            run.font.bold = False
            run.font.color.rgb = COLOR_TEXT_DARK
            
            y_pos += Inches(0.65)
        
        # Перенос на следующий слайд если не помещается
        if y_pos > Inches(7.0):
            break
    
    return slide

def create_two_column_slide(prs, section_num, section_title, left_title, left_content, right_title, right_content):
    """Слайд с двумя колонками"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Светлый фон
    add_background(slide, COLOR_BG_LIGHT)
    
    # Бирюзовый акцент слева
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.4), prs.slide_height)
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_ACCENT
    accent_bar.line.fill.background()
    
    # Заголовок секции
    header_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.4), Inches(11.5), Inches(1))
    header_frame = header_box.text_frame
    header_frame.clear()
    header_para = header_frame.paragraphs[0]
    header_para.text = f"{section_num} / {section_title}"
    
    run = header_para.runs[0]
    run.font.name = FONT_FALLBACK
    run.font.size = HEADING_SIZE
    run.font.bold = True
    run.font.color.rgb = COLOR_TEXT_DARK
    
    # Разделитель
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.3), Inches(11.5), Inches(0.05))
    divider.fill.solid()
    divider.fill.fore_color.rgb = COLOR_ACCENT
    divider.line.fill.background()
    
    # Левая колонка - заголовок
    left_title_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.5), Inches(0.5))
    left_title_frame = left_title_box.text_frame
    left_title_frame.clear()
    left_title_para = left_title_frame.paragraphs[0]
    left_title_para.text = left_title
    
    run = left_title_para.runs[0]
    run.font.name = FONT_FALLBACK
    run.font.size = SUBHEADING_SIZE
    run.font.bold = True
    run.font.color.rgb = COLOR_TEXT_DARK
    
    # Левая колонка - контент с плашкой
    left_placeholder = add_placeholder_box(slide, Inches(1.0), Inches(2.3), Inches(5.5), Inches(4.5))
    left_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.45), Inches(5.1), Inches(4.2))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    left_frame.clear()
    for i, text in enumerate(left_content):
        if i == 0:
            para = left_frame.paragraphs[0]
        else:
            para = left_frame.add_paragraph()
        para.text = text
        run = para.runs[0]
        run.font.name = FONT_FALLBACK
        run.font.size = BODY_SIZE
        run.font.bold = False
        run.font.color.rgb = COLOR_TEXT_DARK
        para.space_after = Pt(8)
    
    # Правая колонка - заголовок
    right_title_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.5), Inches(0.5))
    right_title_frame = right_title_box.text_frame
    right_title_frame.clear()
    right_title_para = right_title_frame.paragraphs[0]
    right_title_para.text = right_title
    
    run = right_title_para.runs[0]
    run.font.name = FONT_FALLBACK
    run.font.size = SUBHEADING_SIZE
    run.font.bold = True
    run.font.color.rgb = COLOR_TEXT_DARK
    
    # Правая колонка - контент с плашкой
    right_placeholder = add_placeholder_box(slide, Inches(7.0), Inches(2.3), Inches(5.5), Inches(4.5))
    right_box = slide.shapes.add_textbox(Inches(7.2), Inches(2.45), Inches(5.1), Inches(4.2))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    right_frame.clear()
    for i, text in enumerate(right_content):
        if i == 0:
            para = right_frame.paragraphs[0]
        else:
            para = right_frame.add_paragraph()
        para.text = text
        run = para.runs[0]
        run.font.name = FONT_FALLBACK
        run.font.size = BODY_SIZE
        run.font.bold = False
        run.font.color.rgb = COLOR_TEXT_DARK
        para.space_after = Pt(8)
    
    return slide

# === СОЗДАНИЕ ВСЕХ СЛАЙДОВ ИЗ drop.docx ===

# Слайд 1: Титульный
slide1 = create_title_slide(
    prs,
    "Культура.Онлайн",
    [
        "Цифровая платформа для учреждений культуры: современный сайт, афиша и онлайн-запись - без разработчика и долгого внедрения",
        "",
        "[Имя основателя] | [Email / сайт] | Стадия: [MVP / пилотные клиенты]"
    ]
)

# Слайд 2: Проблема - суть
slide2 = create_section_slide(
    prs,
    "01",
    "ПРОБЛЕМА — Актуальность, востребованность",
    [
        {'type': 'subheading', 'text': 'В чём суть проблемы:'},
        {'type': 'body', 'text': 'В России ~16 000 учреждений культуры. Большинство из них не существуют в интернете - нет современного сайта, нет афиши, нет онлайн-билетов. Молодёжь ищет в яндексе куда пойти, ничего не находит и идёт в ТЦ или кино.'},
        {'type': 'subheading', 'text': 'Цена проблемы:'},
        {'type': 'body', 'text': 'Учреждения теряют десятки посетителей ежедневно. При среднем билете 400 руб. - это 2–5 млн руб./год упущенной выручки с одного места. А если умножить на 1000? 16000?'},
    ]
)

# Слайд 3: Проблема - почему сейчас
slide3 = create_section_slide(
    prs,
    "01 (продолжение)",
    "ПРОБЛЕМА — Почему решение актуально именно сейчас",
    [
        {'type': 'subheading', 'text': 'Почему решение актуально именно сейчас:'},
        {'type': 'body', 'text': 'Государство активно вкладывает в культуру: Пушкинская карта дала молодёжи деньги на поход в музей, нацпроект «Культура» финансирует ремонты и оборудование. А цифровой инфраструктуры нет. Учреждения физически обновляются, а способы записи к ним как в 2009.'},
    ]
)

# Слайд 4: Решение
slide4 = create_section_slide(
    prs,
    "02",
    "РЕШЕНИЕ — Предлагаемое решение проблемы",
    [
        {'type': 'subheading', 'text': 'Ваше решение:'},
        {'type': 'body', 'text': 'Культура.Онлайн - готовая отраслевая платформа для учреждений культуры. Не разработка с нуля под каждого, а единая система с продуманной архитектурой, которая адаптируется под конкретный музей или дом культуры и запускается за 7–14 дней.'},
        {'type': 'body', 'text': 'Учреждение получает современный сайт, афишу с онлайн-записью и продажей билетов, интеграцию с Пушкинской картой - и панель управления, где сотрудник добавляет новое мероприятие за 3 минуты без участия программиста.'},
        {'type': 'body', 'text': 'Технику берём на себя: хостинг, обновления, поддержка. Учреждение просто работает с аудиторией.'},
    ]
)

# Слайд 5: Продукт - что входит
slide5 = create_section_slide(
    prs,
    "03",
    "ПРОДУКТ",
    [
        {'type': 'body', 'text': 'Архитектура разработана один раз. Каждый новый клиент - адаптация под учреждение, а не разработка с нуля. Это и есть масштаб и отсутствие границ.'},
        {'type': 'subheading', 'text': 'Что входит:'},
        {'type': 'bullet', 'text': 'Современный сайт с готовой структурой под задачи культурного учреждения'},
        {'type': 'bullet', 'text': 'Встроенный модуль афиши: сотрудник добавляет мероприятие сам - как пост в соцсети'},
        {'type': 'bullet', 'text': 'Онлайн-запись/продажа билетов, в том числе по Пушкинской карте'},
        {'type': 'bullet', 'text': 'Панель управления без единой строки кода'},
        {'type': 'bullet', 'text': 'Техническое сопровождение и обновления функций'},
    ]
)

# Слайд 6: Продукт - ключевые модули
slide6 = create_section_slide(
    prs,
    "03 (продолжение)",
    "ПРОДУКТ — Ключевые модули платформы",
    [
        {'type': 'subheading', 'text': 'Ключевые модули платформы:'},
        {'type': 'bullet', 'text': 'Афиша и календарь событий'},
        {'type': 'bullet', 'text': 'Онлайн-бронирование и оплата билетов'},
        {'type': 'bullet', 'text': 'Интеграция с Пушкинской картой'},
        {'type': 'bullet', 'text': 'Новости и публикации'},
        {'type': 'bullet', 'text': 'Адаптивный дизайн'},
        {'type': 'bullet', 'text': 'Аналитика посещаемости'},
        {'type': 'small', 'text': '[Здесь рекомендуется разместить скриншоты панели управления или пример готового сайта]'},
    ]
)

# Слайд 7: Бизнес-модель - монетизация
slide7 = create_section_slide(
    prs,
    "06",
    "БИЗНЕС-МОДЕЛЬ — Как проект зарабатывает деньги?",
    [
        {'type': 'subheading', 'text': 'Монетизация:'},
        {'type': 'body', 'text': 'Модель подписки (ежемесячный платёж за пользование) + разовая плата за внедрение и адаптацию.'},
        {'type': 'bullet', 'text': 'Внедрение и настройка (+оплата за первый месяц): 29.000 рублей'},
        {'type': 'bullet', 'text': 'Ежемесячная подписка: 9.000 рублей/месяц (техподдержка, обновления, хостинг, инструкции по работе) либо добавить возможность оплаты 50.000 рублей/год (выгода 53%).'},
        {'type': 'small', 'text': '! = это меньше, чем разовый запрос разработчику - а сюда уже входит хостинг, обновления, техподдержка и инструкции по работе'},
        {'type': 'bullet', 'text': 'Дополнительные опции по запросу: интеграция с кассой, мультиязычность, 3Д-туры'},
    ]
)

# Слайд 8: Бизнес-модель - клиенты и партнёры
slide8 = create_section_slide(
    prs,
    "06 (продолжение)",
    "БИЗНЕС-МОДЕЛЬ — Клиенты и партнёры",
    [
        {'type': 'subheading', 'text': 'Кто клиенты:'},
        {'type': 'body', 'text': 'Приоритетные сегменты - государственные и муниципальные музеи, галереи и выставочные залы (300–500), дома культуры. Именно музеи - наиболее платёжеспособный сегмент с государственным финансированием.'},
        {'type': 'subheading', 'text': 'Как узнают:'},
        {'type': 'bullet', 'text': 'Прямые переговоры с руководством учреждений (письма по email, телефонные звонки, личные встречи)'},
        {'type': 'bullet', 'text': 'Участие в отраслевых форумах и культурных событиях'},
        {'type': 'bullet', 'text': 'Партнёрство с региональными министерствами культуры'},
        {'type': 'bullet', 'text': 'Рекомендации внутри профессионального сообщества'},
        {'type': 'subheading', 'text': 'Кто партнёры:'},
        {'type': 'body', 'text': 'Потенциальные партнёры - платёжные системы для интеграции билетной части, операторы Пушкинской карты, региональные агентства по развитию культуры.'},
    ]
)

# Слайд 9: Рынок - TAM/SAM/SOM
slide9 = create_section_slide(
    prs,
    "07",
    "РЫНОК — TAM / SAM / SOM",
    [
        {'type': 'subheading', 'text': 'TAM (весь рынок):'},
        {'type': 'body', 'text': 'Учреждения культуры в России - это свыше 16 000 объектов: ~3 900 государственных музеев (Госкаталог, 2025), тысячи домов культуры (только в 2024 отремонтировано 1 562 из них), сотни галерей, заповедников и арт-пространств. Каждое из них потенциально нуждается в цифровом присутствии.'},
        {'type': 'small', 'text': 'Расчёт TAM: 15 000 учреждений × 50 000 руб./год (подписка) = ~750 млн руб./год'},
        {'type': 'subheading', 'text': 'SAM (достижимый рынок):'},
        {'type': 'body', 'text': 'Учреждения с реальным бюджетом на цифру и потребностью в продукте - прежде всего государственные и муниципальные музеи, галереи, крупные дома культуры. Оценочно - 4 000–5 000 объектов.'},
        {'type': 'small', 'text': 'Расчёт SAM: 4 000 учреждений × 50 000 руб./год = ~200 млн руб./год'},
        {'type': 'subheading', 'text': 'SOM (реалистичная доля за 1–3 года):'},
        {'type': 'body', 'text': 'При фокусе на двух-трех регионах и последовательном масштабировании - 200-300 подключённых учреждений к концу третьего года.'},
        {'type': 'small', 'text': 'Расчёт SOM: 200 учреждений × 50 000 руб./год = ~10 млн руб./год'},
        {'type': 'small', 'text': 'Примечание: все цифры приведены на основе данных Госкаталога музейного фонда РФ (ноябрь 2025), Парламентской газеты (данные Минкульта, 2024), реестра частных музеев (ArtInvestment, 2020). Расчёт рыночных объёмов - оценочный, на основе открытых данных.'},
    ]
)

# Слайд 10: Конкуренция
slide10 = create_section_slide(
    prs,
    "08",
    "КОНКУРЕНЦИЯ — Существующие альтернативы",
    [
        {'type': 'body', 'text': 'Главное преимущество: мы не универсальный конструктор и не агентство. Мы продукт, который уже знает, что такое объект культуры, Пушкинская карта и экскурсионное расписание - и не требует объяснять это заново каждому клиенту.'},
    ]
)

# Слайд 11: Результаты
slide11 = create_section_slide(
    prs,
    "09",
    "РЕЗУЛЬТАТЫ — Достигнутые результаты",
    [
        {'type': 'body', 'text': 'Стадия: разработан прототип платформы - спроектированы основные страницы, блоки, структура админ-панели. Остаётся дизайн, вёрстка, сложные технические интеграции. Именно поэтому требуется привлечение финансирования.'},
        {'type': 'subheading', 'text': 'Подтверждённый интерес:'},
        {'type': 'body', 'text': 'Проведено 11 глубинных интервью с представителями учреждений городов: Екатеринбург, Верхняя Пышма, Среднеуральск. Более половины готовы попробовать решение при понятной цене и условиях.'},
        {'type': 'subheading', 'text': 'Рынок подтверждён:'},
        {'type': 'body', 'text': 'Проанализировано 10+ сегментов, свыше ≈10 000 потенциальных клиентов по России.'},
        {'type': 'subheading', 'text': 'Свердловская область - стартовый регион:'},
        {'type': 'body', 'text': 'Более 2 000 организаций культуры, из них свыше 120 государственных и муниципальных музеев и около 900 культурно-досуговых учреждений - и это только один регион.'},
    ]
)

# Слайд 12: Стратегия развития
slide12 = create_section_slide(
    prs,
    "08",
    "СТРАТЕГИЯ РАЗВИТИЯ — Как проект будет масштабироваться?",
    [
        {'type': 'bullet', 'text': 'Год 1. Пилотный запуск в одном регионе, имея минимально рабочий продукт. Подключение первых 10 учреждений. Доработка продукта по обратной связи. Формирование кейсов по сегментам (для лучшего понимания ценности продукта).'},
        {'type': 'bullet', 'text': 'Год 2. Масштабирование на 3-5 регионов. Партнёрства с региональными министерствами культуры. Подключение 50-100 учреждений. Первая стабильная подписная выручка.'},
        {'type': 'bullet', 'text': 'Год 3. Федеральный охват. Продуктовое расширение (аналитика посещаемости, CRM для учреждений, интеграция с государственными системами). 200-300 активных клиентов.'},
        {'type': 'bullet', 'text': 'Год 4–5. Платформенная модель: учреждения сами рекомендуют продукт коллегам. Выход на самоокупаемость и рост без пропорционального увеличения команды - за счёт единой архитектуры.'},
        {'type': 'body', 'text': 'Масштабируемость заложена в архитектуре: добавление нового клиента не требует создания нового продукта. Это означает рост выручки без пропорционального роста затрат.'},
    ]
)

# Слайд 13: Финансы
slide13 = create_section_slide(
    prs,
    "12",
    "ФИНАНСЫ — Прогноз на 1–5 лет",
    [
        {'type': 'body', 'text': 'Расчёт построен на базовом сценарии: средняя подписка 9 000 руб./мес. на клиента + разовое внедрение 29 000 руб.'},
    ]
)

# Слайд 14: Запрос инвестиций - на что
slide14 = create_section_slide(
    prs,
    "13",
    "ЗАПРОС — На что нужны инвестиции?",
    [
        {'type': 'subheading', 'text': 'Запрос: 1.000.000 рублей'},
        {'type': 'subheading', 'text': 'Направления использования:'},
        {'type': 'bullet', 'text': 'Продукт (45% = 450.000) - доведение платформы до логического завершения: дизайн, верстка, модули редактирования, интеграции, мобильное приложение для администратора учреждения, усиление технической части.'},
        {'type': 'bullet', 'text': 'Продажи и развитие (55% = 550.000) - найм менеджера по работе с клиентами, участие в отраслевых мероприятиях, пилотные партнёрства с региональными учреждениями, рассылки email-писем для руководителей, реклама в профильных каналах.'},
    ]
)

# Слайд 15: Запрос инвестиций - почему грант
slide15 = create_section_slide(
    prs,
    "13 (продолжение)",
    "ЗАПРОС — Почему грант, а не кредит?",
    [
        {'type': 'body', 'text': 'Продукт структурно готов - есть прототип, архитектура, логика. Не хватает дизайна, вёрстки и финального запуска. Грант закрывает именно этот этап: довести до рабочего продукта и выйти на первых 50 клиентов.'},
        {'type': 'body', 'text': 'Социальный эффект: больше учреждений культуры онлайн - больше людей до них доходит. Простая механика с измеримым результатом.'},
    ]
)

# Слайд 16: Команда
slide16 = create_two_column_slide(
    prs,
    "11",
    "КОМАНДА — Почему именно вы?",
    "Семён - руководитель проекта",
    [
        "В сфере запуска интернет-проектов с 2019 года.",
        "Руководил командой 6+ человек.",
        "Понимает продукт целиком: как работает, как интегрируется с внешними сервисами.",
        "Контролирует все этапы: от идеи до запуска."
    ],
    "Алексей - архитектура и логика продукта",
    [
        "Спроектировал структуру более 50 сайтов и платформ, среди которых smartranking.ru - федеральный агрегатор рейтингов.",
        "Понимает, где пользователь теряется, а где всё понятно.",
        "Отвечает за логику работы платформы."
    ]
)

# Сохраняем презентацию
output_file = 'prezentaciya_kultura_online_final.pptx'
prs.save(output_file)
print(f"✅ Презентация успешно создана: {output_file}")
print(f"📊 Всего слайдов: {len(prs.slides)}")
print(f"📁 Размер файла: {len(open(output_file, 'rb').read()) / 1024:.1f} KB")
