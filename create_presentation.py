from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Создаем презентацию
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 формат
prs.slide_height = Inches(7.5)

# Цветовая палитра (современный стиль)
COLOR_PRIMARY = RGBColor(20, 24, 50)     # Очень темный синий/черный
COLOR_ACCENT = RGBColor(255, 107, 53)    # Оранжевый акцент
COLOR_TEXT_DARK = RGBColor(30, 30, 30)   # Темно-серый текст
COLOR_TEXT_LIGHT = RGBColor(80, 80, 80)  # Светлее текст
COLOR_BG_LIGHT = RGBColor(250, 250, 252) # Светлый фон
COLOR_WHITE = RGBColor(255, 255, 255)

# Размеры шрифтов
TITLE_SIZE = Pt(54)
SUBTITLE_SIZE = Pt(24)
HEADING_SIZE = Pt(36)
SUBHEADING_SIZE = Pt(24)
BODY_SIZE = Pt(18)
SMALL_SIZE = Pt(14)

def add_background(slide, color=COLOR_WHITE):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg

def set_font_for_paragraph(paragraph, size=BODY_SIZE, bold=False, color=COLOR_TEXT_DARK, font_name='Arial'):
    paragraph.clear()
    run = paragraph.add_run()
    run.font.name = font_name
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    return run

def set_font_for_run(run, size=BODY_SIZE, bold=False, color=COLOR_TEXT_DARK, font_name='Arial'):
    run.font.name = font_name
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    return run

def create_title_slide(prs, title_text, subtitle_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_background(slide, COLOR_PRIMARY)
    
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), prs.slide_height)
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_ACCENT
    accent_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title_text
    title_para.alignment = PP_ALIGN.LEFT
    
    set_font_for_paragraph(title_para, size=TITLE_SIZE, bold=True, color=COLOR_WHITE, font_name='Arial')
    
    sub_y = Inches(4.2)
    for line in subtitle_lines:
        sub_box = slide.shapes.add_textbox(Inches(1), sub_y, Inches(11), Inches(0.5))
        sub_frame = sub_box.text_frame
        sub_para = sub_frame.paragraphs[0]
        sub_para.text = line
        sub_para.alignment = PP_ALIGN.LEFT
        
        set_font_for_paragraph(sub_para, size=SUBTITLE_SIZE, bold=False, color=COLOR_WHITE, font_name='Arial')
        sub_y += Inches(0.4)
    
    return slide

def create_section_slide(prs, section_num, section_title, content_blocks):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_background(slide, COLOR_BG_LIGHT)
    
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), prs.slide_height)
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_ACCENT
    accent_bar.line.fill.background()
    
    header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11.5), Inches(1))
    header_frame = header_box.text_frame
    header_para = header_frame.paragraphs[0]
    header_para.text = f"{section_num} / {section_title}"
    
    set_font_for_paragraph(header_para, size=HEADING_SIZE, bold=True, color=COLOR_PRIMARY, font_name='Arial')
    
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(11.5), Inches(0.03))
    divider.fill.solid()
    divider.fill.fore_color.rgb = COLOR_ACCENT
    divider.line.fill.background()
    
    y_pos = Inches(1.5)
    for block in content_blocks:
        block_type = block.get('type', 'body')
        text = block.get('text', '')
        
        if block_type == 'subheading':
            box = slide.shapes.add_textbox(Inches(0.8), y_pos, Inches(11.5), Inches(0.5))
            frame = box.text_frame
            frame.clear()
            para = frame.paragraphs[0]
            para.text = text
            set_font_for_paragraph(para, size=SUBHEADING_SIZE, bold=True, color=COLOR_PRIMARY, font_name='Arial')
            y_pos += Inches(0.6)
        elif block_type == 'body':
            box = slide.shapes.add_textbox(Inches(0.8), y_pos, Inches(11.5), Inches(0.6))
            frame = box.text_frame
            frame.word_wrap = True
            frame.clear()
            para = frame.paragraphs[0]
            para.text = text
            set_font_for_paragraph(para, size=BODY_SIZE, bold=False, color=COLOR_TEXT_DARK, font_name='Arial')
            y_pos += Inches(0.55)
        elif block_type == 'bullet':
            box = slide.shapes.add_textbox(Inches(1.0), y_pos, Inches(11.3), Inches(0.5))
            frame = box.text_frame
            frame.word_wrap = True
            frame.clear()
            para = frame.paragraphs[0]
            para.text = '• ' + text
            set_font_for_paragraph(para, size=BODY_SIZE, bold=False, color=COLOR_TEXT_LIGHT, font_name='Arial')
            y_pos += Inches(0.5)
        elif block_type == 'small':
            box = slide.shapes.add_textbox(Inches(0.8), y_pos, Inches(11.5), Inches(0.4))
            frame = box.text_frame
            frame.word_wrap = True
            frame.clear()
            para = frame.paragraphs[0]
            para.text = text
            set_font_for_paragraph(para, size=SMALL_SIZE, bold=False, color=COLOR_TEXT_LIGHT, font_name='Arial')
            y_pos += Inches(0.4)
        
        if y_pos > Inches(7.0):
            break
    
    return slide

def create_two_column_slide(prs, section_num, section_title, left_title, left_content, right_title, right_content):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_background(slide, COLOR_BG_LIGHT)
    
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), prs.slide_height)
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_ACCENT
    accent_bar.line.fill.background()
    
    header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11.5), Inches(1))
    header_frame = header_box.text_frame
    header_para = header_frame.paragraphs[0]
    header_para.text = f"{section_num} / {section_title}"
    set_font_for_paragraph(header_para, size=HEADING_SIZE, bold=True, color=COLOR_PRIMARY, font_name='Arial')
    
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(11.5), Inches(0.03))
    divider.fill.solid()
    divider.fill.fore_color.rgb = COLOR_ACCENT
    divider.line.fill.background()
    
    left_title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5))
    left_title_frame = left_title_box.text_frame
    left_title_para = left_title_frame.paragraphs[0]
    left_title_para.text = left_title
    set_font_for_paragraph(left_title_para, size=SUBHEADING_SIZE, bold=True, color=COLOR_PRIMARY, font_name='Arial')
    
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(5.5), Inches(5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    left_frame.clear()
    for i, text in enumerate(left_content):
        if i == 0:
            para = left_frame.paragraphs[0]
        else:
            para = left_frame.add_paragraph()
        para.text = text
        set_font_for_paragraph(para, size=BODY_SIZE, bold=False, color=COLOR_TEXT_LIGHT, font_name='Arial')
        para.space_after = Pt(8)
    
    right_title_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.5))
    right_title_frame = right_title_box.text_frame
    right_title_para = right_title_frame.paragraphs[0]
    right_title_para.text = right_title
    set_font_for_paragraph(right_title_para, size=SUBHEADING_SIZE, bold=True, color=COLOR_PRIMARY, font_name='Arial')
    
    right_box = slide.shapes.add_textbox(Inches(7.0), Inches(2.1), Inches(5.5), Inches(5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    right_frame.clear()
    for i, text in enumerate(right_content):
        if i == 0:
            para = right_frame.paragraphs[0]
        else:
            para = right_frame.add_paragraph()
        para.text = text
        set_font_for_paragraph(para, size=BODY_SIZE, bold=False, color=COLOR_TEXT_LIGHT, font_name='Arial')
        para.space_after = Pt(8)
    
    return slide

# === СОЗДАНИЕ ВСЕХ СЛАЙДОВ ===

slide1 = create_title_slide(
    prs,
    "Культура.Онлайн",
    [
        "Цифровая платформа для учреждений культуры: современный сайт, афиша и онлайн-запись - без разработчика и долгого внедрения",
        "",
        "[Имя основателя] | [Email / сайт] | Стадия: [MVP / пилотные клиенты]"
    ]
)

slide2 = create_section_slide(
    prs,
    "01",
    "ПРОБЛЕМА — Актуальность, востребованность",
    [
        {'type': 'subheading', 'text': 'В чём суть проблемы:'},
        {'type': 'body', 'text': 'В России ~16 000 учреждений культуры. Большинство из них не существуют в интернете - нет современного сайта, нет афиши, нет онлайн-билетов. Молодёжь ищет в яндексе куда пойти, ничего не находит и идёт в ТЦ или кино.'},
        {'type': 'subheading', 'text': 'Цена проблемы:'},
        {'type': 'body', 'text': 'Учреждения теряют десятки посетителей ежедневно. При среднем билете 400 руб. - это 2–5 млн руб./год упущенной выручки с одного места. А если умножить на 1000? 16000?'},
    ]
)

slide3 = create_section_slide(
    prs,
    "01 (продолжение)",
    "ПРОБЛЕМА — Почему сейчас",
    [
        {'type': 'subheading', 'text': 'Почему решение актуально именно сейчас:'},
        {'type': 'body', 'text': 'Государство активно вкладывает в культуру: Пушкинская карта дала молодёжи деньги на поход в музей, нацпроект «Культура» финансирует ремонты и оборудование. А цифровой инфраструктуры нет. Учреждения физически обновляются, а способы записи к ним как в 2009.'},
    ]
)

slide4 = create_section_slide(
    prs,
    "02",
    "РЕШЕНИЕ — Предлагаемое решение проблемы",
    [
        {'type': 'subheading', 'text': 'Ваше решение:'},
        {'type': 'body', 'text': 'Культура.Онлайн - готовая отраслевая платформа для учреждений культуры. Не разработка с нуля под каждого, а единая система с продуманной архитектурой, которая адаптируется под конкретный музей или дом культуры и запускается за 7–14 дней.'},
        {'type': 'body', 'text': 'Учреждение получает современный сайт, афишу с онлайн-записью и продажей билетов, интеграцию с Пушкинской картой - и панель управления, где сотрудник добавляет новое мероприятие за 3 минуты без участия программиста.'},
        {'type': 'body', 'text': 'Технику берём на себя: хостинг, обновления, поддержка. Учреждение просто работает с аудиторией.'},
    ]
)

slide5 = create_section_slide(
    prs,
    "03",
    "ПРОДУКТ",
    [
        {'type': 'body', 'text': 'Архитектура разработана один раз. Каждый новый клиент - адаптация под учреждение, а не разработка с нуля. Это и есть масштаб и отсутствие границ.'},
        {'type': 'subheading', 'text': 'Что входит:'},
        {'type': 'bullet', 'text': 'Современный сайт с готовой структурой под задачи культурного учреждения'},
        {'type': 'bullet', 'text': 'Встроенный модуль афиши: сотрудник добавляет мероприятие сам - как пост в соцсети'},
        {'type': 'bullet', 'text': 'Онлайн-запись/продажа билетов, в том числе по Пушкинской карте'},
        {'type': 'bullet', 'text': 'Панель управления без единой строки кода'},
        {'type': 'bullet', 'text': 'Техническое сопровождение и обновления функций'},
    ]
)

slide6 = create_section_slide(
    prs,
    "03 (продолжение)",
    "ПРОДУКТ — Ключевые модули",
    [
        {'type': 'subheading', 'text': 'Ключевые модули платформы:'},
        {'type': 'bullet', 'text': 'Афиша и календарь событий'},
        {'type': 'bullet', 'text': 'Онлайн-бронирование и оплата билетов'},
        {'type': 'bullet', 'text': 'Интеграция с Пушкинской картой'},
        {'type': 'bullet', 'text': 'Новости и публикации'},
        {'type': 'bullet', 'text': 'Адаптивный дизайн'},
        {'type': 'bullet', 'text': 'Аналитика посещаемости'},
        {'type': 'small', 'text': '[Здесь рекомендуется разместить скриншоты панели управления или пример готового сайта]'},
    ]
)

slide7 = create_section_slide(
    prs,
    "06",
    "БИЗНЕС-МОДЕЛЬ — Как проект зарабатывает деньги?",
    [
        {'type': 'subheading', 'text': 'Монетизация:'},
        {'type': 'body', 'text': 'Модель подписки (ежемесячный платёж за пользование) + разовая плата за внедрение и адаптацию.'},
        {'type': 'bullet', 'text': 'Внедрение и настройка (+оплата за первый месяц): 29.000 рублей'},
        {'type': 'bullet', 'text': 'Ежемесячная подписка: 9.000 рублей/месяц (техподдержка, обновления, хостинг, инструкции по работе) либо добавить возможность оплаты 50.000 рублей/год (выгода 53%).'},
        {'type': 'small', 'text': '! = это меньше, чем разовый запрос разработчику - а сюда уже входит хостинг, обновления, техподдержка и инструкции по работе'},
        {'type': 'bullet', 'text': 'Дополнительные опции по запросу: интеграция с кассой, мультиязычность, 3Д-туры'},
    ]
)

slide8 = create_section_slide(
    prs,
    "06 (продолжение)",
    "БИЗНЕС-МОДЕЛЬ — Клиенты и партнёры",
    [
        {'type': 'subheading', 'text': 'Кто клиенты:'},
        {'type': 'body', 'text': 'Приоритетные сегменты - государственные и муниципальные музеи, галереи и выставочные залы (300–500), дома культуры. Именно музеи - наиболее платёжеспособный сегмент с государственным финансированием.'},
        {'type': 'subheading', 'text': 'Как узнают:'},
        {'type': 'bullet', 'text': 'Прямые переговоры с руководством учреждений (письма по email, телефонные звонки, личные встречи)'},
        {'type': 'bullet', 'text': 'Участие в отраслевых форумах и культурных событиях'},
        {'type': 'bullet', 'text': 'Партнёрство с региональными министерствами культуры'},
        {'type': 'bullet', 'text': 'Рекомендации внутри профессионального сообщества'},
        {'type': 'subheading', 'text': 'Кто партнёры:'},
        {'type': 'body', 'text': 'Потенциальные партнёры - платёжные системы для интеграции билетной части, операторы Пушкинской карты, региональные агентства по развитию культуры.'},
    ]
)

slide9 = create_section_slide(
    prs,
    "07",
    "РЫНОК — TAM / SAM / SOM",
    [
        {'type': 'subheading', 'text': 'TAM (весь рынок):'},
        {'type': 'body', 'text': 'Учреждения культуры в России - это свыше 16 000 объектов: ~3 900 государственных музеев (Госкаталог, 2025), тысячи домов культуры (только в 2024 отремонтировано 1 562 из них), сотни галерей, заповедников и арт-пространств. Каждое из них потенциально нуждается в цифровом присутствии.'},
        {'type': 'small', 'text': 'Расчёт TAM: 15 000 учреждений × 50 000 руб./год (подписка) = ~750 млн руб./год'},
        {'type': 'subheading', 'text': 'SAM (достижимый рынок):'},
        {'type': 'body', 'text': 'Учреждения с реальным бюджетом на цифру и потребностью в продукте - прежде всего государственные и муниципальные музеи, галереи, крупные дома культуры. Оценочно - 4 000–5 000 объектов.'},
        {'type': 'small', 'text': 'Расчёт SAM: 4 000 учреждений × 50 000 руб./год = ~200 млн руб./год'},
        {'type': 'subheading', 'text': 'SOM (реалистичная доля за 1–3 года):'},
        {'type': 'body', 'text': 'При фокусе на двух-трех регионах и последовательном масштабировании - 200-300 подключённых учреждений к концу третьего года.'},
        {'type': 'small', 'text': 'Расчёт SOM: 200 учреждений × 50 000 руб./год = ~10 млн руб./год'},
        {'type': 'small', 'text': 'Примечание: все цифры приведены на основе данных Госкаталога музейного фонда РФ (ноябрь 2025), Парламентской газеты (данные Минкульта, 2024), реестра частных музеев (ArtInvestment, 2020). Расчёт рыночных объёмов - оценочный, на основе открытых данных.'},
    ]
)

slide10 = create_section_slide(
    prs,
    "08",
    "КОНКУРЕНЦИЯ — Существующие альтернативы",
    [
        {'type': 'body', 'text': 'Главное преимущество: мы не универсальный конструктор и не агентство. Мы продукт, который уже знает, что такое объект культуры, Пушкинская карта и экскурсионное расписание - и не требует объяснять это заново каждому клиенту.'},
    ]
)

slide11 = create_section_slide(
    prs,
    "09",
    "РЕЗУЛЬТАТЫ — Достигнутые результаты",
    [
        {'type': 'subheading', 'text': 'Стадия:'},
        {'type': 'body', 'text': 'разработан прототип платформы - спроектированы основные страницы, блоки, структура админ-панели. Остаётся дизайн, вёрстка, сложные технические интеграции. Именно поэтому требуется привлечение средств.'},
        {'type': 'subheading', 'text': 'Подтверждённый интерес:'},
        {'type': 'body', 'text': 'Проведено 11 глубинных интервью с представителями учреждений городов: Екатеринбург, Верхняя Пышма, Среднеуральск. Более половины готовы попробовать решение при понятной цене и простом внедрении.'},
        {'type': 'subheading', 'text': 'Рынок подтверждён:'},
        {'type': 'body', 'text': 'Проанализировано 10+ сегментов, свыше ≈10 000 потенциальных клиентов по России.'},
        {'type': 'subheading', 'text': 'Свердловская область - стартовый регион:'},
        {'type': 'body', 'text': 'Более 2 000 организаций культуры, из них свыше 120 государственных и муниципальных музеев и около 900 культурно-досуговых учреждений - и это только один регион из 89.'},
    ]
)

slide12 = create_section_slide(
    prs,
    "08",
    "СТРАТЕГИЯ РАЗВИТИЯ — Как проект будет масштабироваться?",
    [
        {'type': 'subheading', 'text': 'Год 1. Пилотный запуск в одном регионе, имея минимально рабочий продукт.'},
        {'type': 'body', 'text': 'Подключение первых 10 учреждений. Доработка продукта по обратной связи. Формирование кейсов по сегментам (для лучшего понимания).'},
        {'type': 'subheading', 'text': 'Год 2. Масштабирование на 3-5 регионов.'},
        {'type': 'body', 'text': 'Партнёрства с региональными министерствами культуры. Подключение 50-100 учреждений. Первая стабильная подписная выручка.'},
        {'type': 'subheading', 'text': 'Год 3. Федеральный охват.'},
        {'type': 'body', 'text': 'Продуктовое расширение (аналитика посещаемости, CRM для учреждений, интеграция с государственными системами). 200-300 активных клиентов.'},
        {'type': 'subheading', 'text': 'Год 4–5. Платформенная модель:'},
        {'type': 'body', 'text': 'учреждения сами рекомендуют продукт коллегам. Выход на самоокупаемость и рост без пропорционального увеличения команды - за счёт единой архитектуры.'},
        {'type': 'body', 'text': 'Масштабируемость заложена в архитектуре: добавление нового клиента не требует создания нового продукта. Это означает рост выручки без пропорционального роста затрат.'},
    ]
)

slide13 = create_section_slide(
    prs,
    "12",
    "ФИНАНСЫ — Прогноз на 1–5 лет",
    [
        {'type': 'body', 'text': 'Расчёт построен на базовом сценарии: средняя подписка 9 000 руб./мес. на клиента + разовое внедрение 29 000 руб.'},
    ]
)

slide14 = create_section_slide(
    prs,
    "13",
    "ЗАПРОС — На что нужны инвестиции?",
    [
        {'type': 'subheading', 'text': 'Запрос: 1.000.000 рублей'},
        {'type': 'subheading', 'text': 'Направления использования:'},
        {'type': 'subheading', 'text': 'Продукт (45% = 450.000)'},
        {'type': 'body', 'text': 'доведение платформы до логического завершения: дизайн, верстка, модули редактирования, интеграции, мобильное приложение для администратора учреждения, усиление технической части для обеспечения надёжности платформы при росте числа клиентов.'},
        {'type': 'subheading', 'text': 'Продажи и развитие (55% = 550.000)'},
        {'type': 'body', 'text': 'найм менеджера по работе с клиентами, участие в отраслевых мероприятиях, пилотные партнёрства с региональными учреждениями, рассылки email-писем для руководителей, звонки с рассказом о продукте, запуск рекламы Яндекс Директ, разработка посадочной страницы с рассказом о продукте, сбор обратной связи.'},
    ]
)

slide15 = create_section_slide(
    prs,
    "13 (продолжение)",
    "ЗАПРОС — Почему грант",
    [
        {'type': 'subheading', 'text': 'Почему грант, а не кредит:'},
        {'type': 'body', 'text': 'Продукт структурно готов - есть прототип, архитектура, логика. Не хватает дизайна, вёрстки и финального запуска. Грант закрывает именно этот этап: довести до рабочего продукта и выйти на первых 50 клиентов без долговой нагрузки.'},
        {'type': 'body', 'text': 'Социальный эффект: больше учреждений культуры онлайн - больше людей до них доходит. Простая механика с измеримым результатом.'},
    ]
)

slide16 = create_two_column_slide(
    prs,
    "11",
    "КОМАНДА — Почему именно вы?",
    "Семён - руководитель проекта",
    [
        "В сфере запуска интернет-проектов с 2019 года. Руководил командой 6+ человек. Понимает продукт целиком: как работает, как интегрируется с внешними сервисами. Контролирует проект от идеи до запуска, генерирует решения, когда старые не работают. Отвечает за то, чтобы продукт не застрял между «придумали» и «запустили»."
    ],
    "Алексей - архитектура и логика продукта",
    [
        "Спроектировал структуру более 50 сайтов и платформ, среди которых smartranking.ru - федеральный агрегатор рейтингов. Понимает, где пользователь теряется, а где конвертируется в заявку. Строит продукт так, чтобы посетитель сайта музея не просто зашёл, а дошёл до покупки билета."
    ]
)

# Сохраняем презентацию
output_path = '/workspace/prezentaciya_kultura_online_full.pptx'
prs.save(output_path)
print(f"✅ Презентация успешно сохранена: {output_path}")
print(f"📊 Количество слайдов: {len(prs.slides)}")
import os
print(f"📁 Размер файла: {round(os.path.getsize(output_path) / 1024)} KB")
